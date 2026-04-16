"""
Génère la présentation Loxia.pptx
Usage: python3 generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Palette de couleurs ────────────────────────────────────────────────────
BG_DARK        = RGBColor(0x0D, 0x1B, 0x2A)   # fond principal
BG_CARD        = RGBColor(0x16, 0x28, 0x3C)   # fond des cartes
BG_CARD2       = RGBColor(0x1E, 0x35, 0x4D)   # cartes secondaires
ACCENT_BLUE    = RGBColor(0x38, 0x9B, 0xFF)   # bleu principal
ACCENT_PURPLE  = RGBColor(0x7C, 0x3A, 0xED)   # violet accent
ACCENT_GREEN   = RGBColor(0x10, 0xB9, 0x81)   # vert succès
ACCENT_ORANGE  = RGBColor(0xF5, 0x9E, 0x0B)   # orange attention
ACCENT_RED     = RGBColor(0xEF, 0x44, 0x44)   # rouge
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_300       = RGBColor(0xCB, 0xD5, 0xE1)
GRAY_400       = RGBColor(0x94, 0xA3, 0xB8)
GRAY_500       = RGBColor(0x64, 0x74, 0x8B)

# ─── Dimensions (widescreen 16:9) ───────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # layout vide

# ─── Utilitaires ────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txb

def add_paragraph(tf, text, font_size=14, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return p

def set_bg(slide, color=BG_DARK):
    bg = add_rect(slide, 0, 0, W, H, color)
    bg.name = "bg"

def add_accent_bar(slide, color=ACCENT_BLUE, width=Inches(0.06)):
    add_rect(slide, 0, 0, width, H, color)

def add_slide_header(slide, title, subtitle=None, color=ACCENT_BLUE):
    """Bandeau titre en haut + barre colorée à gauche"""
    add_accent_bar(slide, color)
    # Ligne séparatrice sous le titre
    add_rect(slide, Inches(0.15), Inches(1.05), W - Inches(0.3), Pt(2), color)
    add_text(slide, title,
             x=Inches(0.3), y=Inches(0.12),
             w=W - Inches(0.6), h=Inches(0.75),
             font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 x=Inches(0.3), y=Inches(0.85),
                 w=W - Inches(0.6), h=Inches(0.28),
                 font_size=14, color=GRAY_400)

def card(slide, x, y, w, h, bg=BG_CARD, radius=None):
    return add_rect(slide, x, y, w, h, bg)

def bullet_block(slide, items, x, y, w, h, title=None, bg=BG_CARD,
                 title_color=ACCENT_BLUE, item_size=13.5, title_size=15):
    c = card(slide, x, y, w, h, bg)
    offset = Inches(0.18)
    ty = y + offset
    if title:
        add_text(slide, title, x + offset, ty, w - offset*2,
                 Inches(0.32), font_size=title_size, bold=True, color=title_color)
        ty += Inches(0.36)
    for item in items:
        add_text(slide, item, x + offset, ty, w - offset*2,
                 Inches(0.26), font_size=item_size, color=GRAY_300)
        ty += Inches(0.27)
    return c


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide, BG_DARK)

# Grande zone décorative droite
add_rect(slide, W - Inches(4.8), 0, Inches(4.8), H, BG_CARD)
add_rect(slide, W - Inches(4.8), 0, Pt(2), H, ACCENT_BLUE)

# Barre horizontale déco
add_rect(slide, 0, Inches(3.6), W - Inches(4.8), Pt(2), ACCENT_BLUE)

# Badge "Projet pédagogique"
badge = add_rect(slide, Inches(0.55), Inches(1.6), Inches(2.8), Inches(0.4), ACCENT_PURPLE)
add_text(slide, "PROJET PÉDAGOGIQUE · ÉCOLE D'INGÉNIEUR",
         Inches(0.6), Inches(1.62), Inches(2.7), Inches(0.36),
         font_size=9, bold=True, color=WHITE)

# Titre principal
add_text(slide, "Loxia",
         Inches(0.55), Inches(2.1), Inches(6.5), Inches(1.4),
         font_size=92, bold=True, color=WHITE)

# Sous-titre
add_text(slide, "Plateforme de location immobilière\nen architecture microservices",
         Inches(0.55), Inches(3.7), Inches(6.8), Inches(1.2),
         font_size=22, color=GRAY_300)

# Ligne déco sous-titre
add_rect(slide, Inches(0.55), Inches(3.65), Inches(1.2), Pt(3), ACCENT_GREEN)

# Infos droites
info_items = [
    ("🎓", "Architecture Logicielle", ACCENT_BLUE),
    ("⚙️",  "Java 21 · Spring Boot · React", WHITE),
    ("🐳", "Docker Compose · Microservices", WHITE),
    ("👥", "6 étudiants ingénieurs", ACCENT_GREEN),
]
iy = Inches(1.5)
for icon, label, col in info_items:
    add_text(slide, icon + "  " + label,
             W - Inches(4.5), iy, Inches(4.2), Inches(0.45),
             font_size=15, color=col)
    iy += Inches(0.52)

# Année
add_text(slide, "2025 – 2026",
         Inches(0.55), Inches(6.8), Inches(3), Inches(0.4),
         font_size=13, color=GRAY_500)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Plan de la présentation", color=ACCENT_BLUE)

sections = [
    ("01", "Contexte & objectifs", ACCENT_BLUE),
    ("02", "Périmètre fonctionnel", ACCENT_GREEN),
    ("03", "Architecture globale", ACCENT_PURPLE),
    ("04", "Stack technique", ACCENT_ORANGE),
    ("05", "Sécurité & authentification", ACCENT_RED),
    ("06", "Communication inter-services", ACCENT_BLUE),
    ("07", "Infrastructure Docker", ACCENT_GREEN),
    ("08", "Décisions architecturales", ACCENT_PURPLE),
    ("09", "Démonstration", ACCENT_ORANGE),
    ("10", "Bilan & enseignements", ACCENT_RED),
]

cols = 2
item_w = Inches(5.9)
item_h = Inches(0.52)
gap_x  = Inches(0.25)
gap_y  = Inches(0.12)
start_x = [Inches(0.3), Inches(6.7)]
sy = Inches(1.3)

for i, (num, label, col) in enumerate(sections):
    col_idx = i % cols
    row_idx = i // cols
    x = start_x[col_idx]
    y = sy + row_idx * (item_h + gap_y)
    bg = card(slide, x, y, item_w, item_h, BG_CARD)
    add_rect(slide, x, y, Inches(0.08), item_h, col)
    add_text(slide, num,
             x + Inches(0.18), y + Inches(0.08), Inches(0.5), item_h - Inches(0.1),
             font_size=18, bold=True, color=col)
    add_text(slide, label,
             x + Inches(0.7), y + Inches(0.1), item_w - Inches(0.8), item_h - Inches(0.1),
             font_size=16, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ÉQUIPE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "L'équipe", "6 étudiants ingénieurs", color=ACCENT_GREEN)

membres = [
    ("Ugo Durand",        "Back-end & DevOps"),
    ("Jonas Obrun",       "Back-end (catalog, rental)"),
    ("Nicolas Ramirez",   "Back-end (auth, gateway)"),
    ("Cyril Grosjean",    "Back-end & intégration"),
    ("Ibrahim Khan",      "Front-end React"),
    ("Tarek El Missiry",  "Front-end & notifications"),
]

card_w = Inches(3.9)
card_h = Inches(1.3)
gap = Inches(0.28)
sx = Inches(0.28)
sy = Inches(1.3)
for i, (name, role) in enumerate(membres):
    row = i // 3
    col = i % 3
    x = sx + col * (card_w + gap)
    y = sy + row * (card_h + gap)
    c = card(slide, x, y, card_w, card_h, BG_CARD2)
    add_rect(slide, x, y, card_w, Pt(3), ACCENT_GREEN)
    # Initiales
    initials = "".join(p[0] for p in name.split())
    add_rect(slide, x + Inches(0.18), y + Inches(0.22), Inches(0.62), Inches(0.62), ACCENT_GREEN)
    add_text(slide, initials,
             x + Inches(0.18), y + Inches(0.22), Inches(0.62), Inches(0.62),
             font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, name,
             x + Inches(0.92), y + Inches(0.22),
             card_w - Inches(1.1), Inches(0.38),
             font_size=15, bold=True, color=WHITE)
    add_text(slide, role,
             x + Inches(0.92), y + Inches(0.62),
             card_w - Inches(1.1), Inches(0.32),
             font_size=12, color=GRAY_400)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CONTEXTE & OBJECTIFS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Contexte & objectifs pédagogiques", color=ACCENT_BLUE)

# Colonne gauche — contexte
add_text(slide, "Le projet",
         Inches(0.3), Inches(1.25), Inches(5.8), Inches(0.4),
         font_size=18, bold=True, color=ACCENT_BLUE)

context_lines = [
    "Loxia est une plateforme web de location de logements",
    "développée dans le cadre d'un cours d'architecture",
    "logicielle. Elle met en relation locataires et propriétaires",
    "via une interface moderne et une API distribuée.",
    "",
    "Périmètre : fonctionnel complet, hors paiement / chat /",
    "upload réel / cartographie.",
]
cy = Inches(1.68)
for line in context_lines:
    add_text(slide, line, Inches(0.3), cy, Inches(5.8), Inches(0.26),
             font_size=13.5, color=GRAY_300)
    cy += Inches(0.265)

# Colonne droite — objectifs
objectives = [
    ("🎯", "Bounded contexts DDD",   "4 microservices métier avec responsabilités cloisonnées"),
    ("🔐", "Auth centralisée JWT",   "Signée par auth-service, validée uniquement à la gateway"),
    ("🔗", "Communication REST",     "Spring RestClient, timeouts, fail-safe & degradation"),
    ("🗄️",  "Isolation des données", "1 base PostgreSQL par service, aucune jointure cross-svc"),
    ("🐳", "Conteneurisation",       "Docker Compose, multi-stage builds, healthchecks"),
    ("⚛️",  "Front découplé",        "React + TypeScript + TanStack Query + Zustand"),
]

oy = Inches(1.25)
for icon, title, desc in objectives:
    c = card(slide, Inches(6.5), oy, Inches(6.5), Inches(0.75), BG_CARD)
    add_rect(slide, Inches(6.5), oy, Pt(3), Inches(0.75), ACCENT_BLUE)
    add_text(slide, icon + "  " + title,
             Inches(6.7), oy + Inches(0.06), Inches(6.0), Inches(0.3),
             font_size=13, bold=True, color=WHITE)
    add_text(slide, desc,
             Inches(6.7), oy + Inches(0.36), Inches(6.0), Inches(0.3),
             font_size=11.5, color=GRAY_400)
    oy += Inches(0.82)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PÉRIMÈTRE FONCTIONNEL
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Périmètre fonctionnel", "Ce que l'application fait réellement", color=ACCENT_GREEN)

features = [
    ("🔍", "Recherche & filtres",       "Filtres ville, type de bien, budget min/max"),
    ("🏠", "Détail d'une annonce",      "Photos, surface, pièces, prix, description, propriétaire"),
    ("📝", "Dépôt de candidature",      "Revenus, situation pro, message personnalisé"),
    ("📊", "Dashboard locataire",       "Suivi des candidatures : En attente / Validé / Refusé"),
    ("✏️",  "Création d'annonce",       "Formulaire complet côté propriétaire"),
    ("🔒", "Règle de verrouillage",     "Annonce non modifiable si candidature PENDING/ACCEPTED"),
    ("📥", "Dashboard propriétaire",    "Réception et traitement des candidatures (accept/reject)"),
    ("🔔", "Notifications in-app",      "Cloche, compteur non lus, marquage lu"),
    ("👤", "Compte unifié",             "Même compte → toggle Locataire / Propriétaire"),
    ("🔑", "Auth complète",             "Inscription, connexion, refresh JWT, déconnexion, profil"),
]

col_w = Inches(6.15)
col_gap = Inches(0.25)
item_h2 = Inches(0.62)
gap2 = Inches(0.1)
sx2 = Inches(0.3)
sy2 = Inches(1.3)

for i, (icon, title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = sx2 + col * (col_w + col_gap)
    y = sy2 + row * (item_h2 + gap2)
    c = card(slide, x, y, col_w, item_h2, BG_CARD)
    add_rect(slide, x, y, Pt(3), item_h2, ACCENT_GREEN)
    add_text(slide, icon + "  " + title,
             x + Inches(0.15), y + Inches(0.07), col_w - Inches(0.2), Inches(0.28),
             font_size=13, bold=True, color=WHITE)
    add_text(slide, desc,
             x + Inches(0.15), y + Inches(0.34), col_w - Inches(0.2), Inches(0.24),
             font_size=11, color=GRAY_400)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE GLOBALE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Architecture globale", "4 microservices + Gateway + Frontend React", color=ACCENT_PURPLE)

# Représentation visuelle de l'architecture
# Frontend
add_rect(slide, Inches(5.2), Inches(1.2), Inches(3.0), Inches(0.72), ACCENT_BLUE)
add_text(slide, "Frontend",
         Inches(5.2), Inches(1.2), Inches(3.0), Inches(0.36),
         font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "React · Vite · TypeScript",
         Inches(5.2), Inches(1.52), Inches(3.0), Inches(0.28),
         font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

# Flèche Frontend → Gateway
add_rect(slide, Inches(6.65), Inches(1.92), Pt(2), Inches(0.45), GRAY_500)
add_text(slide, "HTTPS/REST",
         Inches(6.7), Inches(1.97), Inches(1.2), Inches(0.25),
         font_size=9, color=GRAY_500)

# Gateway
add_rect(slide, Inches(4.6), Inches(2.37), Inches(4.2), Inches(0.82), ACCENT_PURPLE)
add_text(slide, "API Gateway  :8080",
         Inches(4.6), Inches(2.37), Inches(4.2), Inches(0.4),
         font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Spring Cloud Gateway · JWT Validation · CORS · Routing",
         Inches(4.6), Inches(2.72), Inches(4.2), Inches(0.32),
         font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

# Flèches Gateway → services (verticale)
add_rect(slide, Inches(6.65), Inches(3.19), Pt(2), Inches(0.5), GRAY_500)

# Services — ligne horizontale
services_info = [
    ("auth-service\n:8081", "Identité\nJWT · Profil", ACCENT_BLUE,    Inches(0.5)),
    ("catalog-service\n:8082", "Annonces\nCRUD · Recherche", ACCENT_GREEN, Inches(3.5)),
    ("rental-service\n:8083", "Candidatures\nLock rule", ACCENT_ORANGE, Inches(6.5)),
    ("notification-\nservice :8084", "Notifications\nIn-app", ACCENT_RED, Inches(9.5)),
]

svc_y = Inches(3.69)
svc_h = Inches(1.1)
svc_w = Inches(2.8)

for svc_name, svc_desc, svc_col, svc_x in services_info:
    add_rect(slide, svc_x, svc_y, svc_w, svc_h, svc_col)
    add_text(slide, svc_name,
             svc_x + Inches(0.1), svc_y + Inches(0.05), svc_w - Inches(0.2), Inches(0.5),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, svc_desc,
             svc_x + Inches(0.1), svc_y + Inches(0.52), svc_w - Inches(0.2), Inches(0.5),
             font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

# Bases de données
dbs = [
    ("auth_db",         ACCENT_BLUE,   Inches(0.5)),
    ("catalog_db",      ACCENT_GREEN,  Inches(3.5)),
    ("rental_db",       ACCENT_ORANGE, Inches(6.5)),
    ("notification_db", ACCENT_RED,    Inches(9.5)),
]
db_y = Inches(5.35)
for db_name, db_col, db_x in dbs:
    # Flèche verticale
    add_rect(slide, db_x + Inches(1.35), svc_y + svc_h, Pt(2), Inches(0.32), GRAY_500)
    add_rect(slide, db_x, db_y, svc_w, Inches(0.62), BG_CARD2)
    add_rect(slide, db_x, db_y, svc_w, Pt(3), db_col)
    add_text(slide, "🗄  " + db_name,
             db_x + Inches(0.1), db_y + Inches(0.12), svc_w - Inches(0.2), Inches(0.38),
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# PostgreSQL unique
add_rect(slide, Inches(3.3), Inches(6.15), Inches(6.8), Inches(0.42), BG_CARD)
add_rect(slide, Inches(3.3), Inches(6.15), Pt(3), Inches(0.42), GRAY_400)
add_text(slide, "PostgreSQL 16 — conteneur unique  loxia-db  (4 bases isolées)",
         Inches(3.45), Inches(6.2), Inches(6.5), Inches(0.32),
         font_size=11, color=GRAY_300)

# Flèches inter-services
add_text(slide, "◄──────────────────────► REST sync  (lock check / owner enrichment / notifications)",
         Inches(0.5), Inches(4.9), Inches(12.0), Inches(0.3),
         font_size=10, color=GRAY_500, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MICROSERVICES EN DÉTAIL
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Les 4 microservices", "Bounded contexts, endpoints, entités", color=ACCENT_PURPLE)

svcs = [
    {
        "name":     "auth-service  :8081",
        "color":    ACCENT_BLUE,
        "context":  "Identité & authentification",
        "entities": "User, RefreshToken",
        "endpoints": [
            "POST  /api/auth/register",
            "POST  /api/auth/login",
            "POST  /api/auth/refresh",
            "GET   /api/auth/me",
            "PUT   /api/auth/me",
        ],
    },
    {
        "name":     "catalog-service  :8082",
        "color":    ACCENT_GREEN,
        "context":  "Annonces immobilières",
        "entities": "Listing",
        "endpoints": [
            "GET   /api/listings",
            "GET   /api/listings/{id}",
            "GET   /api/listings/mine",
            "POST  /api/listings",
            "PUT   /api/listings/{id}",
            "DELETE /api/listings/{id}",
        ],
    },
    {
        "name":     "rental-service  :8083",
        "color":    ACCENT_ORANGE,
        "context":  "Candidatures & cycle de location",
        "entities": "Application",
        "endpoints": [
            "POST  /api/applications",
            "GET   /api/applications/mine",
            "GET   /api/applications/received",
            "POST  /api/applications/{id}/accept",
            "POST  /api/applications/{id}/reject",
        ],
    },
    {
        "name":     "notification-service  :8084",
        "color":    ACCENT_RED,
        "context":  "Notifications in-app",
        "entities": "Notification",
        "endpoints": [
            "GET   /api/notifications/mine",
            "GET   /api/notifications/unread-count",
            "POST  /api/notifications/{id}/read",
            "POST  /api/notifications/read-all",
        ],
    },
]

card_w2 = Inches(3.05)
card_h2 = Inches(5.6)
gap3 = Inches(0.2)
sx3 = Inches(0.25)
for i, svc in enumerate(svcs):
    x = sx3 + i * (card_w2 + gap3)
    y = Inches(1.25)
    c = card(slide, x, y, card_w2, card_h2, BG_CARD)
    col = svc["color"]
    add_rect(slide, x, y, card_w2, Inches(0.48), col)
    add_text(slide, svc["name"],
             x + Inches(0.1), y + Inches(0.06), card_w2 - Inches(0.2), Inches(0.36),
             font_size=12, bold=True, color=WHITE)
    add_text(slide, svc["context"],
             x + Inches(0.1), y + Inches(0.55), card_w2 - Inches(0.2), Inches(0.28),
             font_size=11, italic=True, color=GRAY_400)
    # Entités
    add_text(slide, "Entités : " + svc["entities"],
             x + Inches(0.1), y + Inches(0.88), card_w2 - Inches(0.2), Inches(0.24),
             font_size=10, color=col)
    add_rect(slide, x + Inches(0.1), y + Inches(1.15), card_w2 - Inches(0.2), Pt(1), GRAY_500)
    add_text(slide, "Endpoints exposés",
             x + Inches(0.1), y + Inches(1.22), card_w2 - Inches(0.2), Inches(0.25),
             font_size=11, bold=True, color=WHITE)
    ey = y + Inches(1.52)
    for ep in svc["endpoints"]:
        add_rect(slide, x + Inches(0.1), ey, card_w2 - Inches(0.2), Inches(0.3), BG_CARD2)
        add_text(slide, ep,
                 x + Inches(0.18), ey + Inches(0.03), card_w2 - Inches(0.36), Inches(0.24),
                 font_size=9.5, color=GRAY_300)
        ey += Inches(0.34)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — STACK TECHNIQUE BACKEND
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Stack technique — Backend", "Java 21 · Spring Boot 3.3", color=ACCENT_ORANGE)

back_items = [
    ("☕", "Java 21 LTS",               "Version LTS courante avec Virtual Threads (Project Loom)"),
    ("🌿", "Spring Boot 3.3.x",         "Framework de référence, Spring MVC, Security 6, Data JPA"),
    ("⚙️",  "Maven + Parent POM",        "Multi-module build, versioning centralisé des dépendances"),
    ("🔐", "JWT HS256 (jjwt)",          "Access token 15 min · Refresh token 7 jours (hashé BCrypt)"),
    ("🗄️",  "Spring Data JPA + Flyway", "ORM Hibernate + migrations versionnées par service"),
    ("🔗", "Spring RestClient",         "HTTP client synchrone pour les appels inter-services (Spring 6.1+)"),
    ("✅", "Jakarta Bean Validation",   "Validation des DTO en entrée (@NotBlank, @Min, @Email…)"),
    ("📖", "springdoc-openapi",         "Swagger UI auto-généré par service (/swagger-ui.html)"),
    ("🧪", "JUnit 5 + Mockito",         "Tests unitaires sur la logique métier critique"),
    ("📌", "Lombok",                    "@Slf4j, @Data, @Builder, @RequiredArgsConstructor"),
]

bx = [Inches(0.3), Inches(6.7)]
by = Inches(1.28)
bh = Inches(0.53)
bg2 = Inches(0.1)
bw = Inches(6.1)
for i, (icon, tech, desc) in enumerate(back_items):
    col = i % 2
    row = i // 2
    x = bx[col]
    y = by + row * (bh + bg2)
    c = card(slide, x, y, bw, bh, BG_CARD)
    add_rect(slide, x, y, Pt(3), bh, ACCENT_ORANGE)
    add_text(slide, icon + "  " + tech,
             x + Inches(0.15), y + Inches(0.04), bw - Inches(0.2), Inches(0.26),
             font_size=13, bold=True, color=WHITE)
    add_text(slide, desc,
             x + Inches(0.15), y + Inches(0.29), bw - Inches(0.2), Inches(0.22),
             font_size=11, color=GRAY_400)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — STACK TECHNIQUE FRONTEND
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Stack technique — Frontend", "React 18 · TypeScript · Tailwind CSS", color=ACCENT_BLUE)

front_items = [
    ("⚛️",  "React 18",           "Composants fonctionnels, hooks, rendu concurrent"),
    ("⚡", "Vite",               "Build ultra-rapide, HMR instantané"),
    ("🔷", "TypeScript",         "Typage strict, interfaces pour tous les DTO"),
    ("🎨", "Tailwind CSS",       "Utility-first CSS, responsive design"),
    ("🧩", "shadcn/ui",          "Composants UI accessibles et customisables"),
    ("🔄", "TanStack Query v5",  "Fetching, cache, mutation — tout côté serveur"),
    ("🐻", "Zustand",            "State global client : auth, rôle UI actif (locataire/propriétaire)"),
    ("📡", "Axios",              "HTTP client avec intercepteur JWT et refresh automatique"),
    ("📋", "React Hook Form",    "Formulaires performants + validation Zod"),
    ("🛣️",  "React Router v6",   "Routing SPA, routes protégées, redirect on auth"),
]

for i, (icon, tech, desc) in enumerate(front_items):
    col = i % 2
    row = i // 2
    x = bx[col]
    y = by + row * (bh + bg2)
    c = card(slide, x, y, bw, bh, BG_CARD)
    add_rect(slide, x, y, Pt(3), bh, ACCENT_BLUE)
    add_text(slide, icon + "  " + tech,
             x + Inches(0.15), y + Inches(0.04), bw - Inches(0.2), Inches(0.26),
             font_size=13, bold=True, color=WHITE)
    add_text(slide, desc,
             x + Inches(0.15), y + Inches(0.29), bw - Inches(0.2), Inches(0.22),
             font_size=11, color=GRAY_400)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SÉCURITÉ & AUTHENTIFICATION
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Sécurité & authentification JWT", "Modèle gateway-centric", color=ACCENT_RED)

# Colonne gauche — texte
points = [
    ("Algorithme",    "HS256, clé partagée via JWT_SECRET (env var)"),
    ("Access token",  "Durée : 15 minutes · Claims : sub, email, fullName"),
    ("Refresh token", "Durée : 7 jours · Stocké hashé SHA-256 en base"),
    ("Rotation",      "Chaque refresh révoque l'ancien token (one-time use)"),
    ("Validation",    "Uniquement à la gateway (GlobalFilter)"),
    ("Propagation",   "X-User-Id, X-User-Email, X-User-FullName → services"),
    ("Services down", "Aucune re-validation → font confiance aux headers"),
    ("Protection",    "/internal/** bloqué explicitement (403) par la gateway"),
]

py_start = Inches(1.3)
for title, desc in points:
    c = card(slide, Inches(0.3), py_start, Inches(5.8), Inches(0.54), BG_CARD)
    add_rect(slide, Inches(0.3), py_start, Pt(3), Inches(0.54), ACCENT_RED)
    add_text(slide, title,
             Inches(0.5), py_start + Inches(0.04), Inches(1.4), Inches(0.24),
             font_size=11, bold=True, color=ACCENT_RED)
    add_text(slide, desc,
             Inches(1.95), py_start + Inches(0.04), Inches(4.0), Inches(0.42),
             font_size=11, color=GRAY_300)
    py_start += Inches(0.58)

# Colonne droite — flux
add_text(slide, "Flux de validation",
         Inches(6.5), Inches(1.25), Inches(6.5), Inches(0.35),
         font_size=15, bold=True, color=WHITE)

flow_steps = [
    (ACCENT_RED,    "1. Client → POST /api/auth/login"),
    (GRAY_500,      "    ↓"),
    (ACCENT_RED,    "2. auth-service → génère access + refresh tokens"),
    (GRAY_500,      "    ↓"),
    (ACCENT_ORANGE, "3. Client stocke les tokens (localStorage)"),
    (GRAY_500,      "    ↓"),
    (ACCENT_ORANGE, "4. Requête authentifiée → Authorization: Bearer <token>"),
    (GRAY_500,      "    ↓"),
    (ACCENT_PURPLE, "5. Gateway → valide la signature HS256"),
    (GRAY_500,      "    ↓"),
    (ACCENT_PURPLE, "6. Gateway → propage X-User-Id, X-User-Email, X-User-FullName"),
    (GRAY_500,      "    ↓"),
    (ACCENT_GREEN,  "7. Service downstream → utilise headers (pas de re-validation)"),
    (GRAY_500,      "    ↓"),
    (ACCENT_BLUE,   "8. Token expiré → Axios interceptor → POST /api/auth/refresh"),
]

fy = Inches(1.7)
for col, text in flow_steps:
    add_text(slide, text, Inches(6.5), fy, Inches(6.5), Inches(0.26),
             font_size=11.5, color=col)
    fy += Inches(0.265)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — COMMUNICATION INTER-SERVICES
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Communication inter-services", "REST synchrone · Spring RestClient · Timeouts", color=ACCENT_PURPLE)

# 3 flux principaux
flows = [
    {
        "from": "catalog-service",
        "to":   "rental-service",
        "col":  ACCENT_ORANGE,
        "endpoint": "GET /internal/applications/listing/{id}/locked",
        "trigger":  "PUT ou DELETE /api/listings/{id}",
        "behavior": "Fail-safe : bloque si rental down (timeout 3s)",
        "why":      "Règle de verrouillage — annonce non modifiable si candidature en cours",
    },
    {
        "from": "catalog-service",
        "to":   "auth-service",
        "col":  ACCENT_BLUE,
        "endpoint": "GET /internal/users/{id}",
        "trigger":  "GET /api/listings/{id}",
        "behavior": "Graceful degradation : ownerName = null si auth down",
        "why":      "Enrichissement du nom du propriétaire dans le détail d'annonce",
    },
    {
        "from": "rental-service",
        "to":   "notification-service",
        "col":  ACCENT_RED,
        "endpoint": "POST /internal/notifications",
        "trigger":  "accept ou reject candidature",
        "behavior": "Graceful degradation : notif perdue, transition réussit quand même",
        "why":      "Informer le locataire de la décision du propriétaire",
    },
]

card_w3 = Inches(4.05)
card_h3 = Inches(2.35)
gap4 = Inches(0.2)
sx4 = Inches(0.3)
sy4 = Inches(1.25)

for i, f in enumerate(flows):
    x = sx4 + i * (card_w3 + gap4)
    c = card(slide, x, sy4, card_w3, card_h3, BG_CARD)
    col = f["col"]
    add_rect(slide, x, sy4, card_w3, Pt(3), col)
    # from → to
    add_text(slide, f["from"],
             x + Inches(0.15), sy4 + Inches(0.1), card_w3 * 0.44, Inches(0.3),
             font_size=12, bold=True, color=col)
    add_text(slide, "───────►",
             x + card_w3 * 0.44, sy4 + Inches(0.12), card_w3 * 0.15, Inches(0.26),
             font_size=12, color=GRAY_500)
    add_text(slide, f["to"],
             x + card_w3 * 0.59, sy4 + Inches(0.1), card_w3 * 0.38, Inches(0.3),
             font_size=12, bold=True, color=col)
    add_rect(slide, x + Inches(0.1), sy4 + Inches(0.45), card_w3 - Inches(0.2), Pt(1), GRAY_500)
    labels = [
        ("Endpoint",  f["endpoint"]),
        ("Déclencheur", f["trigger"]),
        ("Timeout",   "3 secondes"),
        ("Comportement", f["behavior"]),
        ("Objectif",  f["why"]),
    ]
    ly = sy4 + Inches(0.55)
    for lbl, val in labels:
        add_text(slide, lbl + " :",
                 x + Inches(0.15), ly, Inches(1.15), Inches(0.24),
                 font_size=10, bold=True, color=col)
        add_text(slide, val,
                 x + Inches(1.3), ly, card_w3 - Inches(1.45), Inches(0.24),
                 font_size=10, color=GRAY_300)
        ly += Inches(0.3)

# Batch endpoint
c2 = card(slide, Inches(0.3), Inches(3.75), Inches(12.7), Inches(0.65), BG_CARD2)
add_rect(slide, Inches(0.3), Inches(3.75), Pt(3), Inches(0.65), ACCENT_PURPLE)
add_text(slide, "Optimisation batch",
         Inches(0.5), Inches(3.8), Inches(2.0), Inches(0.24),
         font_size=12, bold=True, color=ACCENT_PURPLE)
add_text(slide, "POST /internal/applications/locks  (catalog → rental)",
         Inches(2.55), Inches(3.8), Inches(5.5), Inches(0.24),
         font_size=11, color=WHITE)
add_text(slide, "→ Vérification en une seule requête pour N annonces (liste). Réduit le nombre d'appels sur les pages listant plusieurs annonces.",
         Inches(2.55), Inches(4.06), Inches(10.0), Inches(0.24),
         font_size=11, color=GRAY_400)

# Politique générale
add_rect(slide, Inches(0.3), Inches(4.55), Inches(12.7), Inches(0.55), BG_CARD)
add_text(slide, "Choix architectural : REST synchrone uniquement (pas de Kafka/RabbitMQ).",
         Inches(0.5), Inches(4.6), Inches(7.0), Inches(0.24),
         font_size=12, bold=True, color=WHITE)
add_text(slide, "Trade-off délibéré : simplicité d'infrastructure pour un projet pédagogique. Les flows sont peu nombreux (3 flux) et la latence acceptable.",
         Inches(0.5), Inches(4.86), Inches(12.3), Inches(0.22),
         font_size=11, color=GRAY_400)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — RÈGLE DE VERROUILLAGE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Règle métier : le verrouillage d'annonce", "La règle star du projet", color=ACCENT_ORANGE)

# Définition
c = card(slide, Inches(0.3), Inches(1.25), Inches(12.7), Inches(0.88), BG_CARD2)
add_rect(slide, Inches(0.3), Inches(1.25), Pt(4), Inches(0.88), ACCENT_ORANGE)
add_text(slide, "📋  Définition",
         Inches(0.55), Inches(1.3), Inches(3.0), Inches(0.3),
         font_size=14, bold=True, color=ACCENT_ORANGE)
add_text(slide, "Une annonce ayant au moins une candidature PENDING ou ACCEPTED ne peut être ni modifiée (PUT) ni supprimée (DELETE).",
         Inches(0.55), Inches(1.6), Inches(12.1), Inches(0.4),
         font_size=13, color=WHITE)

# Séquence visuelle
add_text(slide, "Flux — PUT /api/listings/{id} avec annonce verrouillée :",
         Inches(0.3), Inches(2.25), Inches(12.0), Inches(0.3),
         font_size=13, bold=True, color=WHITE)

seq_steps = [
    (ACCENT_BLUE,   "Client",               "PUT /api/listings/{id}"),
    (ACCENT_PURPLE, "Gateway",              "Valide JWT · propage X-User-Id"),
    (ACCENT_GREEN,  "catalog-service",      "Vérifie ownership (listing.ownerId == X-User-Id)"),
    (ACCENT_ORANGE, "→ rental-service",     "GET /internal/applications/listing/{id}/locked"),
    (ACCENT_ORANGE, "← rental-service",     "{ locked: true }   (candidature PENDING détectée)"),
    (ACCENT_RED,    "catalog-service",      "Rejette la modification → 409 Conflict"),
    (ACCENT_RED,    "Client",               "Reçoit 409 · l'UI affiche un message explicite"),
]
sx5 = Inches(0.3)
sy5 = Inches(2.65)
sw5 = Inches(12.7)
sh5 = Inches(0.44)
for step_col, actor, action in seq_steps:
    c = card(slide, sx5, sy5, sw5, sh5, BG_CARD)
    add_rect(slide, sx5, sy5, Pt(3), sh5, step_col)
    add_text(slide, actor,
             sx5 + Inches(0.15), sy5 + Inches(0.08), Inches(2.2), Inches(0.26),
             font_size=11, bold=True, color=step_col)
    add_text(slide, action,
             sx5 + Inches(2.4), sy5 + Inches(0.08), sw5 - Inches(2.6), Inches(0.26),
             font_size=11, color=GRAY_300)
    sy5 += Inches(0.47)

# Note sur fail-safe
c3 = card(slide, Inches(0.3), Inches(6.05), Inches(12.7), Inches(0.56), BG_CARD2)
add_rect(slide, Inches(0.3), Inches(6.05), Pt(4), Inches(0.56), ACCENT_RED)
add_text(slide, "⚠️  Fail-safe",
         Inches(0.55), Inches(6.1), Inches(1.5), Inches(0.26),
         font_size=12, bold=True, color=ACCENT_RED)
add_text(slide, "Si rental-service ne répond pas (timeout 3s) → catalog bloque quand même la modification. Mieux vaut refuser une opération légitime que permettre une modification sur un état inconnu.",
         Inches(2.1), Inches(6.1), Inches(10.7), Inches(0.42),
         font_size=11, color=GRAY_300)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — INFRASTRUCTURE DOCKER
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Infrastructure Docker Compose", "8 conteneurs, 3 ports exposés, healthchecks", color=ACCENT_GREEN)

containers = [
    ("loxia-db",             "postgres:16",              "—",       "Base de données PostgreSQL (4 bases)", GRAY_400),
    ("pgadmin",              "dpage/pgadmin4:8",          "8090",    "Interface admin BD (dev only)",        GRAY_400),
    ("auth-service",         "build local",              "—",       "Port interne 8081 uniquement",         ACCENT_BLUE),
    ("catalog-service",      "build local",              "—",       "Port interne 8082 uniquement",         ACCENT_GREEN),
    ("rental-service",       "build local",              "—",       "Port interne 8083 uniquement",         ACCENT_ORANGE),
    ("notification-service", "build local",              "—",       "Port interne 8084 uniquement",         ACCENT_RED),
    ("gateway",              "build local",              "8080",    "Seul point d'entrée HTTP public",      ACCENT_PURPLE),
    ("frontend",             "build local (Nginx)",      "3000",    "Sert le build React statique",         ACCENT_BLUE),
]

# Header table
hx = Inches(0.3)
hy = Inches(1.25)
cols_w = [Inches(2.4), Inches(2.3), Inches(1.0), Inches(4.5), Inches(2.2)]
col_names = ["Conteneur", "Image", "Port hôte", "Rôle", ""]
add_rect(slide, hx, hy, sum(cols_w), Inches(0.38), ACCENT_GREEN)
cx_off = hx
for cn, cw in zip(col_names, cols_w):
    add_text(slide, cn, cx_off + Inches(0.08), hy + Inches(0.05), cw, Inches(0.28),
             font_size=11, bold=True, color=WHITE)
    cx_off += cw

ry = hy + Inches(0.38)
rh = Inches(0.5)
for cname, image, port, role, col in containers:
    bg3 = BG_CARD if (containers.index((cname, image, port, role, col)) % 2 == 0) else BG_CARD2
    add_rect(slide, hx, ry, sum(cols_w), rh, bg3)
    vals = [cname, image, port, role]
    cx_off2 = hx
    for val, cw, vc in zip(vals, cols_w, [col, GRAY_400, ACCENT_ORANGE if port != "—" else GRAY_500, GRAY_300]):
        add_text(slide, val, cx_off2 + Inches(0.08), ry + Inches(0.12), cw, Inches(0.26),
                 font_size=10.5, color=vc)
        cx_off2 += cw
    ry += rh

# Healthchecks
c4 = card(slide, Inches(0.3), Inches(5.85), Inches(12.7), Inches(0.65), BG_CARD2)
add_rect(slide, Inches(0.3), Inches(5.85), Pt(4), Inches(0.65), ACCENT_GREEN)
add_text(slide, "Ordre de démarrage garanti via  depends_on: condition: service_healthy",
         Inches(0.55), Inches(5.9), Inches(7.0), Inches(0.28),
         font_size=12, bold=True, color=WHITE)
add_text(slide, "loxia-db → {auth, catalog, rental, notification} → gateway → frontend",
         Inches(0.55), Inches(6.15), Inches(12.1), Inches(0.28),
         font_size=12, color=ACCENT_GREEN)

# Multi-stage builds
c5 = card(slide, Inches(0.3), Inches(6.62), Inches(12.7), Inches(0.55), BG_CARD)
add_text(slide, "Multi-stage Dockerfiles  —  Stage 1 : Maven build (JDK 21)  |  Stage 2 : JRE Alpine, user non-root loxia",
         Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.28),
         font_size=11.5, color=GRAY_300)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — DÉCISIONS ARCHITECTURALES (ADRs)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Décisions architecturales clés (ADRs)", "10 décisions documentées", color=ACCENT_PURPLE)

adrs = [
    ("ADR-001", "Java 21 + Spring Boot 3.3",  "Standard de facto en entreprise · LTS · Virtual Threads", ACCENT_BLUE),
    ("ADR-002", "4 microservices métier",       "Découpage DDD · autonomie par membre · périmètre maîtrisable", ACCENT_GREEN),
    ("ADR-003", "REST sync (pas de Kafka)",     "Infrastructure simple · 3 flux seulement · couplage temporel accepté", ACCENT_ORANGE),
    ("ADR-004", "Compte unifié (pas de rôles)", "Un user peut louer ET posséder · auth contextuelle (ownership)", ACCENT_RED),
    ("ADR-005", "JWT validé à la gateway",      "Pas de duplication de secret · isolation réseau Docker sécurise", ACCENT_PURPLE),
    ("ADR-006", "1 conteneur PostgreSQL, 4 DB", "Isolation logique sans complexité opérationnelle de 4 conteneurs", ACCENT_BLUE),
    ("ADR-007", "Spring Cloud Gateway",         "Écosystème Spring cohérent · filtre JWT en Java · pédagogique", ACCENT_GREEN),
    ("ADR-008", "DNS Docker Compose (pas Eureka)", "Instances stables → service discovery inutile pour ce scope", ACCENT_ORANGE),
    ("ADR-009", "Pas de Spring Cloud Config",   "1 environnement → 3 couches application.yml + profil + .env", ACCENT_RED),
    ("ADR-010", "Fail-safe sur la règle de lock", "Bloquer si incertain → préférable à permettre un état corrompu", ACCENT_PURPLE),
]

adr_w = Inches(6.1)
adr_h = Inches(0.52)
adr_gap = Inches(0.1)
ax = [Inches(0.3), Inches(6.7)]
ay = Inches(1.28)

for i, (num, title, reason, col) in enumerate(adrs):
    xc = ax[i % 2]
    yc = ay + (i // 2) * (adr_h + adr_gap)
    c = card(slide, xc, yc, adr_w, adr_h, BG_CARD)
    add_rect(slide, xc, yc, Pt(3), adr_h, col)
    add_text(slide, num,
             xc + Inches(0.12), yc + Inches(0.04), Inches(0.82), Inches(0.24),
             font_size=10, bold=True, color=col)
    add_text(slide, title,
             xc + Inches(0.95), yc + Inches(0.04), adr_w - Inches(1.05), Inches(0.24),
             font_size=12, bold=True, color=WHITE)
    add_text(slide, reason,
             xc + Inches(0.12), yc + Inches(0.28), adr_w - Inches(0.25), Inches(0.22),
             font_size=10, color=GRAY_400)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — DETTES TECHNIQUES & RISQUES
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Dettes techniques & risques connus", "Trade-offs documentés et assumés", color=ACCENT_RED)

risks = [
    ("🔄", "Race condition — règle de lock",
     "Entre la vérification /locked et l'UPDATE, une candidature peut arriver.",
     "Documentée. Solution : verrou applicatif ou Saga pattern.",
     "connue"),
    ("🍪", "Refresh token en localStorage",
     "Vulnérable au XSS (vol de token côté client).",
     "En prod : cookie HttpOnly SameSite=Strict.",
     "connue"),
    ("🔑", "JWT HS256 avec secret partagé",
     "Si le secret fuit, n'importe qui peut forger un token.",
     "En prod : RS256 + JWKS endpoint (clé asymétrique).",
     "connue"),
    ("🔍", "Pas de tracing distribué",
     "Difficile de déboguer les appels en cascade entre services.",
     "En prod : OpenTelemetry + Jaeger.",
     "hors scope"),
    ("📦", "Dérive DTO back/front",
     "Types partagés synchronisés manuellement dans frontend/src/types/.",
     "En prod : génération automatique depuis OpenAPI spec.",
     "acceptable"),
    ("🔗", "Monolithe distribué potentiel",
     "Chaîner 3 appels inter-services sur un seul endpoint.",
     "Max 2 appels par endpoint · endpoints batch pour les listes.",
     "mitigé"),
]

rc_w = Inches(6.05)
rc_h = Inches(0.95)
rc_gap = Inches(0.15)
rx = [Inches(0.3), Inches(6.65)]
ry2 = Inches(1.28)

badge_colors = {
    "connue":     (ACCENT_RED,    "DETTE CONNUE"),
    "hors scope": (GRAY_500,      "HORS SCOPE"),
    "acceptable": (ACCENT_ORANGE, "ACCEPTABLE"),
    "mitigé":     (ACCENT_GREEN,  "MITIGÉ"),
}

for i, (icon, title, problem, solution, status) in enumerate(risks):
    xc = rx[i % 2]
    yc = ry2 + (i // 2) * (rc_h + rc_gap)
    c = card(slide, xc, yc, rc_w, rc_h, BG_CARD)
    sc, sl = badge_colors[status]
    add_rect(slide, xc, yc, Pt(3), rc_h, sc)
    # Badge
    add_rect(slide, xc + rc_w - Inches(1.3), yc + Inches(0.08), Inches(1.2), Inches(0.24), sc)
    add_text(slide, sl,
             xc + rc_w - Inches(1.28), yc + Inches(0.1), Inches(1.18), Inches(0.2),
             font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, icon + "  " + title,
             xc + Inches(0.15), yc + Inches(0.06), rc_w - Inches(1.5), Inches(0.26),
             font_size=12, bold=True, color=WHITE)
    add_text(slide, "Problème : " + problem,
             xc + Inches(0.15), yc + Inches(0.36), rc_w - Inches(0.25), Inches(0.24),
             font_size=10, color=GRAY_300)
    add_text(slide, "Solution : " + solution,
             xc + Inches(0.15), yc + Inches(0.62), rc_w - Inches(0.25), Inches(0.24),
             font_size=10, color=ACCENT_GREEN, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — DÉMONSTRATION / SCREENSHOTS
# ════════════════════════════════════════════════════════════════════════════
MOCKUPS = "/Users/cyrilgrosjean/Developer/1-Learning/Loxia/docs/mockups/"

screens = [
    ("Capture d'écran 2026-04-10 113622.png", "Accueil — Vue Locataire (recherche + filtres)"),
    ("Capture d'écran 2026-04-10 113719.png", "Détail d'une annonce (bouton Postuler)"),
    ("Capture d'écran 2026-04-10 113725.png", "Formulaire de candidature"),
    ("Capture d'écran 2026-04-10 113732.png", "Dashboard Propriétaire — Demandes reçues"),
    ("Capture d'écran 2026-04-10 113703.png", "Création d'une annonce"),
    ("Capture d'écran 2026-04-10 113747.png", "Dashboard Locataire — Mes candidatures"),
]

for fname, caption in screens:
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, BG_DARK)
    add_slide_header(slide, "Démonstration", caption, color=ACCENT_BLUE)
    path = os.path.join(MOCKUPS, fname)
    if os.path.exists(path):
        try:
            img = slide.shapes.add_picture(path, Inches(0.5), Inches(1.3),
                                           width=Inches(12.3))
            # Ne pas dépasser la hauteur
            if img.top + img.height > Inches(7.3):
                ratio = img.width / img.height
                img.height = Inches(7.3) - img.top
                img.width  = int(img.height * ratio)
        except Exception as e:
            add_text(slide, f"[image : {fname}]",
                     Inches(0.5), Inches(2), Inches(12), Inches(1),
                     font_size=16, color=GRAY_400)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — CHIFFRES CLÉS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Loxia en chiffres", "", color=ACCENT_GREEN)

stats = [
    ("5",  "services\n(4 métier + gateway)", ACCENT_BLUE),
    ("19", "endpoints\nexposés",             ACCENT_GREEN),
    ("6",  "endpoints\ninternes",            ACCENT_PURPLE),
    ("4",  "bases de données\nisolées",      ACCENT_ORANGE),
    ("8",  "conteneurs\nDocker",             ACCENT_RED),
    ("10", "ADRs\ndocumentées",              ACCENT_BLUE),
    ("16", "maquettes\nde référence",        ACCENT_GREEN),
    ("6",  "étudiants\ningénieurs",          ACCENT_PURPLE),
]

sw2 = Inches(2.9)
sh2 = Inches(2.2)
sg = Inches(0.35)
sx6 = Inches(0.4)
sy6 = Inches(1.35)
for i, (num, label, col) in enumerate(stats):
    row = i // 4
    col_idx = i % 4
    x = sx6 + col_idx * (sw2 + sg)
    y = sy6 + row * (sh2 + Inches(0.2))
    c = card(slide, x, y, sw2, sh2, BG_CARD)
    add_rect(slide, x, y, sw2, Pt(4), col)
    add_text(slide, num,
             x + Inches(0.1), y + Inches(0.2), sw2 - Inches(0.2), Inches(1.0),
             font_size=68, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, label,
             x + Inches(0.1), y + Inches(1.25), sw2 - Inches(0.2), Inches(0.8),
             font_size=13, color=GRAY_300, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — BILAN & ENSEIGNEMENTS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_slide_header(slide, "Bilan & enseignements", "Ce qu'on a appris, ce qu'on ferait différemment", color=ACCENT_BLUE)

# Colonne gauche — ce qui a bien marché
card(slide, Inches(0.3), Inches(1.25), Inches(5.85), Inches(5.5), BG_CARD)
add_rect(slide, Inches(0.3), Inches(1.25), Pt(3), Inches(5.5), ACCENT_GREEN)
add_text(slide, "✅  Ce qui a bien marché",
         Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.38),
         font_size=15, bold=True, color=ACCENT_GREEN)
good = [
    "Découpage en bounded contexts très naturel une fois\nque chaque service a un propriétaire dans l'équipe",
    "La gateway comme seul validateur JWT simplifie\nconsidérablement les services downstream",
    "Docker Compose avec healthchecks → zero surprise\nlors des intégrations",
    "TanStack Query + Zustand → état front prévisible\net clair sans boilerplate Redux",
    "Flyway → migrations versionnées, aucun conflit de\nschéma entre branches",
    "Le compte unifié évite la gestion de rôles complexe\ntout en couvrant tous les use cases",
]
gy = Inches(1.75)
for g in good:
    add_rect(slide, Inches(0.55), gy, Inches(0.08), Inches(0.28), ACCENT_GREEN)
    add_text(slide, g, Inches(0.75), gy, Inches(5.2), Inches(0.52),
             font_size=11.5, color=GRAY_300)
    gy += Inches(0.62)

# Colonne droite — ce qu'on ferait différemment
card(slide, Inches(6.7), Inches(1.25), Inches(5.85), Inches(5.5), BG_CARD)
add_rect(slide, Inches(6.7), Inches(1.25), Pt(3), Inches(5.5), ACCENT_ORANGE)
add_text(slide, "💡  Ce qu'on ferait différemment",
         Inches(6.9), Inches(1.3), Inches(5.5), Inches(0.38),
         font_size=15, bold=True, color=ACCENT_ORANGE)
diff = [
    "Générer les types TypeScript automatiquement\ndepuis les specs OpenAPI de chaque service",
    "Adopter RS256 (clé asymétrique) pour le JWT\ndès le début plutôt que HS256 partagé",
    "Ajouter OpenTelemetry + Jaeger dès les premiers\nservices — le debug distributed est difficile",
    "Introduire un message broker (RabbitMQ) pour les\nnotifications — découplage naturel",
    "Cookies HttpOnly pour les tokens dès le départ,\npas de localStorage",
    "Un vrai CI/CD (GitHub Actions) avec build\nautonome à chaque PR sur develop",
]
dy = Inches(1.75)
for d in diff:
    add_rect(slide, Inches(6.95), dy, Inches(0.08), Inches(0.28), ACCENT_ORANGE)
    add_text(slide, d, Inches(7.15), dy, Inches(5.2), Inches(0.52),
             font_size=11.5, color=GRAY_300)
    dy += Inches(0.62)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — SLIDE DE CLÔTURE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide, BG_DARK)

add_rect(slide, W - Inches(4.8), 0, Inches(4.8), H, BG_CARD)
add_rect(slide, W - Inches(4.8), 0, Pt(2), H, ACCENT_BLUE)
add_rect(slide, 0, Inches(3.6), W - Inches(4.8), Pt(2), ACCENT_BLUE)

add_text(slide, "Merci !",
         Inches(0.55), Inches(1.8), Inches(6.5), Inches(1.4),
         font_size=88, bold=True, color=WHITE)

add_text(slide, "Des questions ?",
         Inches(0.55), Inches(3.75), Inches(6.5), Inches(0.7),
         font_size=28, color=GRAY_300)

add_rect(slide, Inches(0.55), Inches(3.72), Inches(1.5), Pt(3), ACCENT_GREEN)

info_r = [
    ("📁", "github.com/UgoDurand/Loxia",   ACCENT_BLUE),
    ("📖", "docs/architecture.md",          WHITE),
    ("🐳", "docker compose up --build",     ACCENT_GREEN),
]
iy2 = Inches(1.6)
for icon, val, col in info_r:
    add_text(slide, icon + "  " + val,
             W - Inches(4.5), iy2, Inches(4.3), Inches(0.45),
             font_size=15, color=col)
    iy2 += Inches(0.6)

add_text(slide, "Ugo · Jonas · Nicolas · Cyril · Ibrahim · Tarek",
         Inches(0.55), Inches(6.6), Inches(6.5), Inches(0.4),
         font_size=13, color=GRAY_500)

# ─── Sauvegarde ─────────────────────────────────────────────────────────────
out = "/Users/cyrilgrosjean/Developer/1-Learning/Loxia/Loxia_Presentation.pptx"
prs.save(out)
print(f"✅  Présentation générée : {out}")
print(f"   Nombre de slides : {len(prs.slides)}")
